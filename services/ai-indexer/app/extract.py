from __future__ import annotations

import csv
import io
import os
import subprocess
import tempfile
from pathlib import Path


EXTRACTABLE_TYPES = {"document", "spreadsheet", "presentation"}
EXTRACTABLE_EXTENSIONS = {
    "pdf", "docx", "doc", "txt", "md", "rtf", "csv", "pptx", "ppt", "xlsx", "xls"
}


def should_process(file_name: str, file_type: str) -> bool:
    extension = Path(file_name or "").suffix.lower().lstrip(".")
    return file_type in EXTRACTABLE_TYPES or extension in EXTRACTABLE_EXTENSIONS


def extract_text(data: bytes, file_name: str) -> str:
    extension = Path(file_name or "").suffix.lower().lstrip(".")
    if extension == "pdf":
        return _extract_pdf(data)
    if extension in {"docx", "doc"}:
        return _extract_docx(data, extension)
    if extension in {"txt", "md", "rtf"}:
        return _decode(data)
    if extension == "csv":
        return _extract_csv(data)
    if extension in {"pptx", "ppt"}:
        return _extract_pptx(data, extension)
    if extension in {"xlsx", "xls"}:
        return _extract_xlsx(data, extension)
    return ""


def _decode(data: bytes) -> str:
    return data.decode("utf-8", errors="replace") or data.decode("latin-1", errors="ignore")


def _extract_pdf(data: bytes) -> str:
    from pypdf import PdfReader

    pages: list[str] = []
    for page in PdfReader(io.BytesIO(data)).pages:
        text = page.extract_text() or ""
        if text.strip():
            pages.append(text.strip())
    return "\n\n".join(pages)


def _extract_docx(data: bytes, extension: str) -> str:
    if extension == "doc":
        return _extract_legacy_office(data, ".doc")
    from docx import Document

    document = Document(io.BytesIO(data))
    parts = [paragraph.text.strip() for paragraph in document.paragraphs if paragraph.text.strip()]
    for table in document.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n\n".join(parts)


def _extract_csv(data: bytes) -> str:
    reader = csv.reader(io.StringIO(_decode(data)))
    return "\n".join(", ".join(cell for cell in row) for row in reader if any(cell.strip() for cell in row))


def _extract_pptx(data: bytes, extension: str) -> str:
    if extension == "ppt":
        return _extract_legacy_office(data, ".ppt")
    from pptx import Presentation

    presentation = Presentation(io.BytesIO(data))
    slides: list[str] = []
    for number, slide in enumerate(presentation.slides, 1):
        parts = [f"[Slide {number}]"]
        for shape in slide.shapes:
            text = getattr(shape, "text", "")
            if text and text.strip():
                parts.append(text.strip())
        slides.append("\n".join(parts))
    return "\n\n".join(slides)


def _extract_xlsx(data: bytes, extension: str) -> str:
    if extension == "xls":
        return _extract_legacy_office(data, ".xls")
    import openpyxl

    workbook = openpyxl.load_workbook(io.BytesIO(data), read_only=True, data_only=True)
    sheets: list[str] = []
    for sheet in workbook.worksheets:
        rows: list[str] = []
        for row in sheet.iter_rows(values_only=True):
            values = [str(value) for value in row if value is not None]
            if values:
                rows.append(" | ".join(values))
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
    return "\n\n".join(sheets)


def _extract_legacy_office(data: bytes, suffix: str) -> str:
    """Use the installed LibreOffice converter for binary Office formats."""
    with tempfile.TemporaryDirectory(prefix="kibegi-ai-") as directory:
        source = Path(directory) / f"input{suffix}"
        source.write_bytes(data)
        result = subprocess.run(
            ["libreoffice", "--headless", "--convert-to", "txt:Text", "--outdir", directory, str(source)],
            capture_output=True,
            timeout=60,
            check=False,
        )
        if result.returncode != 0:
            return ""
        converted = Path(directory) / "input.txt"
        return _decode(converted.read_bytes()) if converted.exists() else ""


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 150) -> list[str]:
    text = text.strip()
    if not text:
        return []
    chunks: list[str] = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunk = text[start:end]
        if end < len(text):
            for boundary in ("\n\n", "\n", ". ", "? ", "! "):
                position = chunk.rfind(boundary)
                if position > chunk_size // 2:
                    chunk = chunk[: position + len(boundary)]
                    break
        cleaned = chunk.strip()
        if cleaned:
            chunks.append(cleaned)
        if end >= len(text):
            break
        next_start = start + len(chunk) - overlap
        start = next_start if next_start > start else end
    return chunks
