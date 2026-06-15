"""
Text extraction from uploaded files.
Supports: PDF, DOCX, TXT, CSV, PPTX, XLSX (basic)
"""
import io
import logging
import os

logger = logging.getLogger(__name__)


def extract_text(file_obj, file_name: str) -> str:
    """
    Extract plain text from a file-like object.
    Returns empty string for unsupported types or on error.
    """
    ext = os.path.splitext(file_name)[1].lower().lstrip('.')

    try:
        if ext == 'pdf':
            return _extract_pdf(file_obj)
        elif ext in ('docx', 'doc'):
            return _extract_docx(file_obj)
        elif ext in ('txt', 'md', 'rtf'):
            return _extract_text_plain(file_obj)
        elif ext == 'csv':
            return _extract_csv(file_obj)
        elif ext in ('pptx', 'ppt'):
            return _extract_pptx(file_obj)
        elif ext in ('xlsx', 'xls'):
            return _extract_xlsx(file_obj)
        else:
            return ''
    except Exception as exc:
        logger.warning("Text extraction failed for %s (%s): %s", file_name, ext, exc)
        return ''


def _extract_pdf(file_obj) -> str:
    import pdfplumber
    text_parts = []
    with pdfplumber.open(file_obj) as pdf:
        for page in pdf.pages:
            page_text = page.extract_text()
            if page_text:
                text_parts.append(page_text.strip())
    return '\n\n'.join(text_parts)


def _extract_docx(file_obj) -> str:
    from docx import Document
    doc = Document(file_obj)
    paragraphs = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    # Also grab text from tables
    for table in doc.tables:
        for row in table.rows:
            row_text = ' | '.join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return '\n\n'.join(paragraphs)


def _extract_text_plain(file_obj) -> str:
    raw = file_obj.read()
    try:
        return raw.decode('utf-8')
    except UnicodeDecodeError:
        return raw.decode('latin-1', errors='ignore')


def _extract_csv(file_obj) -> str:
    import csv
    raw = file_obj.read()
    try:
        text = raw.decode('utf-8')
    except UnicodeDecodeError:
        text = raw.decode('latin-1', errors='ignore')
    reader = csv.reader(io.StringIO(text))
    rows = [', '.join(row) for row in reader if any(cell.strip() for cell in row)]
    return '\n'.join(rows)


def _extract_pptx(file_obj) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ''
    prs = Presentation(file_obj)
    slides = []
    for i, slide in enumerate(prs.slides, 1):
        parts = [f"[Slide {i}]"]
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                parts.append(shape.text.strip())
        slides.append('\n'.join(parts))
    return '\n\n'.join(slides)


def _extract_xlsx(file_obj) -> str:
    try:
        import openpyxl
    except ImportError:
        return ''
    wb = openpyxl.load_workbook(file_obj, read_only=True, data_only=True)
    sheets = []
    for sheet in wb.worksheets:
        rows = []
        for row in sheet.iter_rows(values_only=True):
            row_text = ' | '.join(str(v) for v in row if v is not None)
            if row_text.strip():
                rows.append(row_text)
        if rows:
            sheets.append(f"[Sheet: {sheet.title}]\n" + '\n'.join(rows))
    return '\n\n'.join(sheets)


def chunk_text(text: str, chunk_size: int = 1200, overlap: int = 150) -> list[str]:
    """
    Split text into overlapping chunks.
    chunk_size ≈ 300 tokens, overlap ensures context continuity.
    """
    text = text.strip()
    if not text:
        return []

    chunks = []
    start = 0
    text_length = len(text)
    while start < text_length:
        end = min(start + chunk_size, text_length)
        chunk = text[start:end]

        # Try to break at a sentence or paragraph boundary when there is more
        # text after this chunk. For the final chunk, append and stop.
        if end < text_length:
            for boundary in ('\n\n', '\n', '. ', '? ', '! '):
                pos = chunk.rfind(boundary)
                if pos > chunk_size // 2:
                    chunk = chunk[:pos + len(boundary)]
                    break

        cleaned = chunk.strip()
        if cleaned:
            chunks.append(cleaned)

        if end >= text_length:
            break

        next_start = start + len(chunk) - overlap
        if next_start <= start:
            next_start = end
        start = next_start

    return chunks
